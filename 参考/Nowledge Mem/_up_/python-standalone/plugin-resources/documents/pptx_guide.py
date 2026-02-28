#!/usr/bin/env python3
"""
PPTX Generation Knowledge Base
Generated: 2025-11-29
Purpose: Document best practices and patterns for python-pptx usage

Run this script to see comprehensive python-pptx documentation.
The agent should run this when uncertain about PPTX generation.
"""

def print_pptx_knowledge():
    """Print comprehensive PPTX generation knowledge and best practices"""

    knowledge = """
# 🎨 PYTHON-PPTX GENERATION GUIDE

## 1. CORE CONCEPTS & ARCHITECTURE

### What is python-pptx?
- Low-level XML manipulation library for PowerPoint (.pptx) files
- Based on OpenXML (OOXML) format
- NOT a design tool - it's a content insertion tool
- Perfect for programmatic slide generation and data population

### Key Classes
```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Main objects:
# - Presentation: Root container for slides
# - Slide: Individual slide in presentation
# - Shape: Any element on a slide (text box, image, table)
# - TextFrame: Container for text within shapes
# - Paragraph: Text block within TextFrame
```

## 2. DESIGN PHILOSOPHY

### ⚠️ CRITICAL LIMITATION
python-pptx is NOT designed to create "elegant, fancy" presentations by itself.
It's designed to:
✅ Add content to slides
✅ Format text (size, color, bold, italic)
✅ Insert images and tables
❌ Create modern, professional designs

### Design Approach (RECOMMENDED)
**Don't fight the library. Work WITH it:**
1. Create beautiful base template in PowerPoint (human designer)
2. Use python-pptx to populate template with data
3. Result: Professional output without design code

### Alternative: Image-Based Rendering
```python
# Generate charts as high-quality PNG
import matplotlib.pyplot as plt
fig, ax = plt.subplots()
ax.plot(data)
plt.savefig('chart.png', dpi=300, bbox_inches='tight')

# Embed in slide
slide.shapes.add_picture('chart.png', Inches(0.5), Inches(1.5), width=Inches(8))
```

## 3. WORKFLOW PATTERNS

### Pattern 1: From Scratch (Basic)
```python
from pptx import Presentation

# Create new presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "My Presentation"
subtitle.text = "Generated with python-pptx"

# Save
prs.save('output.pptx')
```

### Pattern 2: From Template (PROFESSIONAL)
```python
from pptx import Presentation

# Load existing template
template = Presentation('my_template.pptx')

# Populate template
for slide in template.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            # Replace placeholders with dynamic data
            if '{{title}}' in shape.text:
                shape.text = shape.text.replace('{{title}}', dynamic_title)

template.save('output.pptx')
```

### Pattern 3: Blank Slide Layout (Flexible)
```python
from pptx import Presentation

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # 6 = blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Add shapes manually for complete control
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(1)

# Add text box
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Custom Title"
p.font.size = Pt(44)
p.font.bold = True
```

## 4. TEXT FORMATTING (CRITICAL)

### Font Styling
```python
from pptx.util import Pt
from pptx.dml.color import RGBColor

p = paragraph
p.text = "Formatted Text"

# Font properties
p.font.size = Pt(24)
p.font.bold = True
p.font.italic = True
p.font.color.rgb = RGBColor(31, 78, 121)  # Dark blue
p.font.name = 'Calibri'

# Alignment
from pptx.enum.text import PP_ALIGN
p.alignment = PP_ALIGN.CENTER  # or LEFT, RIGHT, CENTER
```

### Color System (ESSENTIAL)
```python
from pptx.dml.color import RGBColor

# Use RGB tuples for custom colors
TITLE_COLOR = RGBColor(31, 78, 121)      # Dark blue
ACCENT_COLOR = RGBColor(192, 0, 0)       # Red
TEXT_COLOR = RGBColor(89, 89, 89)        # Gray
SUCCESS_COLOR = RGBColor(46, 204, 113)   # Green

# Apply to shapes
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE_COLOR
```

### Paragraph Styling
```python
p = text_frame.paragraphs[0]
p.level = 0  # Indent level (0=no indent, 1=sub-bullet, etc)
p.space_before = Pt(12)
p.space_after = Pt(12)
p.line_spacing = 1.5
```

## 5. IMAGES & CHARTS

### Embedding Images
```python
# Must use absolute paths
slide.shapes.add_picture(
    '/absolute/path/to/image.png',
    Inches(0.5),      # left position
    Inches(1.5),      # top position
    width=Inches(8)   # width (height auto-scales)
)
```

### Chart Generation Best Practice
```python
# Step 1: Generate chart as PNG using Matplotlib
import matplotlib.pyplot as plt
plt.figure(figsize=(8, 5))
plt.bar(categories, values, color='#3498DB')
plt.savefig('chart.png', dpi=150, bbox_inches='tight')
plt.close()

# Step 2: Embed in slide
slide.shapes.add_picture('chart.png', Inches(0.8), Inches(1.2), width=Inches(8.4))

# Benefits:
# ✅ Full control over chart appearance
# ✅ High quality output
# ✅ Can use any charting library
# ✅ Charts are "locked" (can't accidentally edit)
```

## 6. TABLES

### Table Creation
```python
from pptx.util import Inches

rows, cols = 4, 3
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set header row
header_cells = table.rows[0].cells
header_cells[0].text = 'Column 1'
header_cells[1].text = 'Column 2'
header_cells[2].text = 'Column 3'

# Fill data
for i, (col1, col2, col3) in enumerate(data, 1):
    row_cells = table.rows[i].cells
    row_cells[0].text = str(col1)
    row_cells[1].text = str(col2)
    row_cells[2].text = str(col3)
```

### Table Styling
```python
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# Cell text formatting
cell = table.rows[0].cells[0]
cell.text = 'Header'
paragraph = cell.text_frame.paragraphs[0]
paragraph.font.bold = True
paragraph.font.size = Pt(12)
```

## 7. SLIDE LAYOUTS

### Available Layouts
```python
prs = Presentation()

# Common layouts:
# 0 = Title Slide
# 1 = Title and Content (bullets)
# 2 = Title Only
# 3 = Blank
# 4 = Title and Two Content
# 5 = Comparison
# 6 = Blank (no placeholders)

layout_0 = prs.slide_layouts[0]  # Title slide
layout_1 = prs.slide_layouts[1]  # Bullet points
layout_6 = prs.slide_layouts[6]  # Blank - full control
```

## 8. BEST PRACTICES

### DO ✅
- Use template-based approach for professional output
- Generate charts as images separately, then embed
- Use color constants (define once, reuse everywhere)
- Test presentation by opening in PowerPoint
- Use absolute file paths
- Save to temp directory during generation, move when done
- Version control your presentation templates

### DON'T ❌
- Try to make python-pptx create "fancy" designs
- Use relative paths (use absolute with os.path.abspath())
- Forget to close matplotlib figures (plt.close())
- Mix presentation concerns (design vs content)
- Generate charts inline with python-pptx
- Assume all fonts exist on all systems

## 9. COMPLETE EXAMPLE: Report Generator

```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import os

def generate_report(title, data, output_path):
    \"\"\"Generate a professional report presentation.\"\"\"

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 ratio
    prs.slide_height = Inches(7.5)

    # Color constants
    TITLE_COLOR = RGBColor(31, 78, 121)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Data Chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Generate chart
    plt.figure(figsize=(10, 6))
    plt.bar(data.keys(), data.values(), color='#3498DB')
    plt.title('Data Overview')
    chart_path = '/tmp/chart.png'
    plt.savefig(chart_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Add chart to slide
    slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1), width=Inches(10))

    # Save
    prs.save(output_path)
    print(f"Created: {output_path}")

# Usage
data = {"Q1": 100, "Q2": 150, "Q3": 120, "Q4": 200}
generate_report("Quarterly Report", data, "report.pptx")
```

## 10. QUICK REFERENCE

| Task | Code |
|------|------|
| Create presentation | `Presentation()` |
| Load template | `Presentation('template.pptx')` |
| Add slide | `prs.slides.add_slide(layout)` |
| Add text box | `slide.shapes.add_textbox(left, top, width, height)` |
| Add image | `slide.shapes.add_picture(path, left, top, width)` |
| Add table | `slide.shapes.add_table(rows, cols, left, top, width, height)` |
| Save | `prs.save('output.pptx')` |
| Set font size | `paragraph.font.size = Pt(24)` |
| Set color | `shape.fill.fore_color.rgb = RGBColor(255, 0, 0)` |
| Align text | `paragraph.alignment = PP_ALIGN.CENTER` |

---

**Remember:** Templates + python-pptx is the winning formula for professional output! 🚀
"""

    print(knowledge)


if __name__ == '__main__':
    print_pptx_knowledge()

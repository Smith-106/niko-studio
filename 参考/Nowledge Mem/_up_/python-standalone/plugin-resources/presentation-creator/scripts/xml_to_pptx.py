#!/usr/bin/env python3
"""
XML to PPTX Converter

Converts structured presentation XML to actual PowerPoint files.

Usage:
    python xml_to_pptx.py input.xml output.pptx
    python xml_to_pptx.py --stdin output.pptx  # Read XML from stdin

Example:
    python xml_to_pptx.py presentation.xml slides.pptx
"""

import sys
import os
import argparse
from xml.etree import ElementTree as ET
from typing import Optional, Dict, List, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.util import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor as DMLRGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# =============================================================================
# Theme Configuration
# =============================================================================

DEFAULT_THEME = {
    'colors': {
        'primary': '#1F4E79',
        'secondary': '#2E75B6',
        'accent': '#70AD47',
        'text': '#333333',
        'background': '#FFFFFF',
    },
    'fonts': {
        'heading': 'Calibri',
        'body': 'Calibri Light',
    }
}

THEMES = {
    'professional': {
        'colors': {
            'primary': '#1F4E79',
            'secondary': '#2E75B6',
            'accent': '#70AD47',
            'text': '#333333',
            'background': '#FFFFFF',
        },
        'fonts': {'heading': 'Calibri', 'body': 'Calibri Light'},
    },
    'minimal': {
        'colors': {
            'primary': '#2D3436',
            'secondary': '#636E72',
            'accent': '#00B894',
            'text': '#2D3436',
            'background': '#FFFFFF',
        },
        'fonts': {'heading': 'Helvetica Neue', 'body': 'Helvetica Neue'},
    },
    'creative': {
        'colors': {
            'primary': '#6C5CE7',
            'secondary': '#A29BFE',
            'accent': '#FD79A8',
            'text': '#2D3436',
            'background': '#FFFFFF',
        },
        'fonts': {'heading': 'Avenir Next', 'body': 'Avenir'},
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def get_theme_color(theme: dict, color_name: str) -> RGBColor:
    """Get a color from theme by name."""
    color_hex = theme['colors'].get(color_name, theme['colors']['text'])
    return hex_to_rgb(color_hex)


# =============================================================================
# Slide Builders
# =============================================================================

class PresentationBuilder:
    """Builds a PowerPoint presentation from parsed XML."""

    def __init__(self, theme: dict):
        self.theme = theme
        self.prs = Presentation()
        # Set 16:9 aspect ratio by default
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def set_aspect_ratio(self, aspect: str):
        """Set slide dimensions based on aspect ratio."""
        if aspect == '4:3':
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        else:  # Default to 16:9
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, slide_elem: ET.Element):
        """Add a title slide."""
        # Use blank layout for full control
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_elem = slide_elem.find('title')
        if title_elem is not None and title_elem.text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5),
                Inches(12.333), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_elem.text.strip()
            p.font.size = Pt(int(title_elem.get('font-size', '44')))
            p.font.bold = True
            p.font.name = self.theme['fonts']['heading']

            color_name = title_elem.get('color', 'primary')
            if color_name.startswith('#'):
                p.font.color.rgb = hex_to_rgb(color_name)
            else:
                p.font.color.rgb = get_theme_color(self.theme, color_name)

            align = title_elem.get('align', 'center')
            p.alignment = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
            }.get(align, PP_ALIGN.CENTER)

        # Subtitle
        subtitle_elem = slide_elem.find('subtitle')
        if subtitle_elem is not None and subtitle_elem.text:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle_elem.text.strip()
            p.font.size = Pt(int(subtitle_elem.get('font-size', '24')))
            p.font.name = self.theme['fonts']['body']
            p.font.color.rgb = get_theme_color(self.theme, 'text')
            p.alignment = PP_ALIGN.CENTER

        self._add_notes(slide, slide_elem)
        return slide

    def add_content_slide(self, slide_elem: ET.Element):
        """Add a content slide with bullets."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_elem = slide_elem.find('title')
        if title_elem is not None and title_elem.text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4),
                Inches(12.333), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_elem.text.strip()
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = self.theme['fonts']['heading']
            p.font.color.rgb = get_theme_color(self.theme, 'primary')

        # Content with bullets
        content_elem = slide_elem.find('content')
        if content_elem is not None:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(12.333), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            bullets = content_elem.findall('bullet')
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = bullet.text.strip() if bullet.text else ''
                p.font.size = Pt(20)
                p.font.name = self.theme['fonts']['body']
                p.font.color.rgb = get_theme_color(self.theme, 'text')

                level = int(bullet.get('level', '0'))
                p.level = level

                # Add spacing
                p.space_before = Pt(6)
                p.space_after = Pt(6)

        # Table
        table_elem = slide_elem.find('table')
        if table_elem is not None:
            self._add_table(slide, table_elem, Inches(0.5), Inches(1.5))

        self._add_notes(slide, slide_elem)
        return slide

    def add_two_column_slide(self, slide_elem: ET.Element):
        """Add a two-column comparison slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_elem = slide_elem.find('title')
        if title_elem is not None and title_elem.text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4),
                Inches(12.333), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_elem.text.strip()
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = self.theme['fonts']['heading']
            p.font.color.rgb = get_theme_color(self.theme, 'primary')

        # Left column
        left_col = slide_elem.find("column[@position='left']")
        if left_col is not None:
            self._add_column(slide, left_col, Inches(0.5), Inches(1.5), Inches(5.9))

        # Right column
        right_col = slide_elem.find("column[@position='right']")
        if right_col is not None:
            self._add_column(slide, right_col, Inches(6.9), Inches(1.5), Inches(5.9))

        self._add_notes(slide, slide_elem)
        return slide

    def add_quote_slide(self, slide_elem: ET.Element):
        """Add a quote slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Quote
        quote_elem = slide_elem.find('quote')
        if quote_elem is not None and quote_elem.text:
            quote_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2),
                Inches(10.333), Inches(3)
            )
            tf = quote_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{quote_elem.text.strip()}"'
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.name = self.theme['fonts']['body']
            p.font.color.rgb = get_theme_color(self.theme, 'text')
            p.alignment = PP_ALIGN.CENTER

        # Attribution
        attr_elem = slide_elem.find('attribution')
        if attr_elem is not None and attr_elem.text:
            attr_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5.2),
                Inches(10.333), Inches(0.8)
            )
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {attr_elem.text.strip()}"
            p.font.size = Pt(18)
            p.font.name = self.theme['fonts']['body']
            p.font.color.rgb = get_theme_color(self.theme, 'secondary')
            p.alignment = PP_ALIGN.CENTER

        self._add_notes(slide, slide_elem)
        return slide

    def add_image_slide(self, slide_elem: ET.Element):
        """Add an image slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Image (placeholder text if file doesn't exist)
        image_elem = slide_elem.find('image')
        if image_elem is not None:
            src = image_elem.get('src', '')
            if src and os.path.exists(src):
                # Add actual image
                slide.shapes.add_picture(
                    src,
                    Inches(0.5), Inches(1),
                    width=Inches(12.333)
                )
            else:
                # Placeholder
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(1),
                    Inches(12.333), Inches(5.5)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
                placeholder.line.color.rgb = RGBColor(200, 200, 200)

                # Add text to placeholder
                tf = placeholder.text_frame
                p = tf.paragraphs[0]
                p.text = f"[Image: {src}]" if src else "[Image placeholder]"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(128, 128, 128)
                p.alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title (overlay style)
        title_elem = slide_elem.find('title')
        if title_elem is not None and title_elem.text:
            position = title_elem.get('position', 'top')
            y_pos = Inches(0.3) if position == 'top' else Inches(6.2)

            title_box = slide.shapes.add_textbox(
                Inches(0.5), y_pos,
                Inches(12.333), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_elem.text.strip()
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = self.theme['fonts']['heading']
            p.font.color.rgb = get_theme_color(self.theme, 'primary')

        # Caption
        caption_elem = slide_elem.find('caption')
        if caption_elem is not None and caption_elem.text:
            cap_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8),
                Inches(12.333), Inches(0.5)
            )
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption_elem.text.strip()
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = get_theme_color(self.theme, 'secondary')
            p.alignment = PP_ALIGN.CENTER

        self._add_notes(slide, slide_elem)
        return slide

    def _add_column(self, slide, col_elem: ET.Element, left: Inches, top: Inches, width: Inches):
        """Add a column to a slide."""
        col_box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
        tf = col_box.text_frame
        tf.word_wrap = True

        # Heading
        heading_elem = col_elem.find('heading')
        if heading_elem is not None and heading_elem.text:
            p = tf.paragraphs[0]
            p.text = heading_elem.text.strip()
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = self.theme['fonts']['heading']
            p.font.color.rgb = get_theme_color(self.theme, 'secondary')
            p.space_after = Pt(12)

        # Bullets
        bullets = col_elem.findall('bullet')
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet.text.strip() if bullet.text else ''
            p.font.size = Pt(18)
            p.font.name = self.theme['fonts']['body']
            p.font.color.rgb = get_theme_color(self.theme, 'text')
            p.level = int(bullet.get('level', '0'))
            p.space_before = Pt(4)
            p.space_after = Pt(4)

        # Table in column
        table_elem = col_elem.find('table')
        if table_elem is not None:
            self._add_table(slide, table_elem, left, top + Inches(0.8), width)

    def _add_table(self, slide, table_elem: ET.Element, left: Inches, top: Inches, width: Inches = None):
        """Add a table to a slide."""
        header = table_elem.find('header')
        rows = table_elem.findall('row')

        if header is None:
            return

        header_cells = header.findall('cell')
        num_cols = len(header_cells)
        num_rows = 1 + len(rows)

        if num_cols == 0:
            return

        table_width = width if width else Inches(10)
        table_height = Inches(0.4 * num_rows)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            left, top,
            table_width, table_height
        )
        table = table_shape.table

        # Header row
        for col_idx, cell_elem in enumerate(header_cells):
            cell = table.cell(0, col_idx)
            cell.text = cell_elem.text.strip() if cell_elem.text else ''
            # Style header
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = get_theme_color(self.theme, 'primary')

        # Data rows
        for row_idx, row_elem in enumerate(rows, 1):
            cells = row_elem.findall('cell')
            for col_idx, cell_elem in enumerate(cells):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_elem.text.strip() if cell_elem.text else ''
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(11)
                    p.font.color.rgb = get_theme_color(self.theme, 'text')

    def _add_notes(self, slide, slide_elem: ET.Element):
        """Add speaker notes to a slide."""
        notes_elem = slide_elem.find('notes')
        if notes_elem is not None and notes_elem.text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_elem.text.strip()

    def save(self, output_path: str):
        """Save the presentation to a file."""
        self.prs.save(output_path)


# =============================================================================
# XML Parser
# =============================================================================

def parse_presentation_xml(xml_content: str) -> Tuple[ET.Element, dict]:
    """Parse presentation XML and extract theme."""
    root = ET.fromstring(xml_content)

    # Get base theme
    theme_name = root.get('theme', 'professional')
    theme = THEMES.get(theme_name, THEMES['professional']).copy()
    theme['colors'] = theme['colors'].copy()
    theme['fonts'] = theme['fonts'].copy()

    # Override with custom theme if present
    theme_elem = root.find('theme')
    if theme_elem is not None:
        colors_elem = theme_elem.find('colors')
        if colors_elem is not None:
            for key in ['primary', 'secondary', 'accent', 'text', 'background']:
                if colors_elem.get(key):
                    theme['colors'][key] = colors_elem.get(key)

        fonts_elem = theme_elem.find('fonts')
        if fonts_elem is not None:
            if fonts_elem.get('heading'):
                theme['fonts']['heading'] = fonts_elem.get('heading')
            if fonts_elem.get('body'):
                theme['fonts']['body'] = fonts_elem.get('body')

    return root, theme


def convert_xml_to_pptx(xml_content: str, output_path: str) -> None:
    """Convert XML content to a PowerPoint file."""
    root, theme = parse_presentation_xml(xml_content)

    builder = PresentationBuilder(theme)

    # Set aspect ratio
    aspect = root.get('aspect', '16:9')
    builder.set_aspect_ratio(aspect)

    # Process each slide
    for slide_elem in root.findall('slide'):
        layout = slide_elem.get('layout', 'content')

        if layout == 'title':
            builder.add_title_slide(slide_elem)
        elif layout == 'two-column':
            builder.add_two_column_slide(slide_elem)
        elif layout == 'quote':
            builder.add_quote_slide(slide_elem)
        elif layout == 'image':
            builder.add_image_slide(slide_elem)
        else:  # content, blank, etc.
            builder.add_content_slide(slide_elem)

    builder.save(output_path)
    print(f"Created: {output_path}")
    print(f"Slides: {len(root.findall('slide'))}")


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description='Convert presentation XML to PowerPoint',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python xml_to_pptx.py presentation.xml output.pptx
  cat presentation.xml | python xml_to_pptx.py --stdin output.pptx
  echo '<presentation>...</presentation>' | python xml_to_pptx.py --stdin slides.pptx
        """
    )
    parser.add_argument('input', nargs='?', help='Input XML file path')
    parser.add_argument('output', help='Output PPTX file path')
    parser.add_argument('--stdin', action='store_true', help='Read XML from stdin')

    args = parser.parse_args()

    if args.stdin:
        xml_content = sys.stdin.read()
    elif args.input:
        with open(args.input, 'r', encoding='utf-8') as f:
            xml_content = f.read()
    else:
        parser.error("Either provide input file or use --stdin")

    try:
        convert_xml_to_pptx(xml_content, args.output)
    except ET.ParseError as e:
        print(f"Error: Invalid XML - {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
